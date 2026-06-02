"""Build docs/slides/chores.pptx — a single-slide deck about the nightly
chore-reset Cloud Function.

Run from the repo root:

    python docs/slides/build_chores_slide.py

Requires: python-pptx (pip install python-pptx).
"""
from __future__ import annotations

import math
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


NAVY = RGBColor(0x1F, 0x2A, 0x44)
INK = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PEOPLE = [
    ("A", 90, RGBColor(0x4C, 0x8B, 0xF5)),
    ("B", 210, RGBColor(0xF5, 0x8B, 0x4C)),
    ("C", 330, RGBColor(0x55, 0xB4, 0x6E)),
]


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    """Straight connector with a triangle arrowhead at the tail end."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = MUTED
    conn.line.width = Pt(2.25)
    ln = conn.line._get_or_add_ln()
    tail_end = etree.SubElement(ln, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("h", "med")


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title --------------------------------------------------------------
    title_tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.1))
    title_tf = title_tb.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = "Chores rotate on a nightly server reset."
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = NAVY

    # Bullets ------------------------------------------------------------
    bullet_tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(7.5), Inches(5))
    btf = bullet_tb.text_frame
    btf.word_wrap = True
    bullets = [
        "One Cloud Function — clients only read.",
        "Anchored day-key math decides who fires today.",
        "Idempotent — catches up missed nights.",
    ]
    for i, text in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = f"•  {text}"
        para.space_after = Pt(18)
        for r in para.runs:
            r.font.size = Pt(24)
            r.font.color.rgb = INK

    # Rotation diagram ---------------------------------------------------
    cx, cy = 10.5, 4.5   # center of diagram, in inches
    ring_r = 1.6
    d = 1.05             # circle diameter

    positions: list[tuple[float, float]] = []
    for label, angle_deg, color in PEOPLE:
        rad = math.radians(angle_deg)
        x = cx + ring_r * math.cos(rad) - d / 2
        y = cy - ring_r * math.sin(rad) - d / 2  # invert: PPT y grows downward
        positions.append((x, y))
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.color.rgb = WHITE
        circ.line.width = Pt(2)
        ctf = circ.text_frame
        ctf.text = label
        for para in ctf.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(32)
                r.font.bold = True
                r.font.color.rgb = WHITE

    # Arrows A -> B -> C -> A
    for i, j in [(0, 1), (1, 2), (2, 0)]:
        xa, ya = positions[i]
        xb, yb = positions[j]
        ca = (xa + d / 2, ya + d / 2)
        cb = (xb + d / 2, yb + d / 2)
        dx, dy = cb[0] - ca[0], cb[1] - ca[1]
        length = math.hypot(dx, dy)
        ux, uy = dx / length, dy / length
        pad = d / 2 + 0.08
        sx, sy = ca[0] + ux * pad, ca[1] + uy * pad
        ex, ey = cb[0] - ux * pad, cb[1] - uy * pad
        add_arrow(slide, sx, sy, ex, ey)

    # Diagram caption
    cap_tb = slide.shapes.add_textbox(
        Inches(cx - 1.5), Inches(cy + ring_r + 0.25), Inches(3), Inches(0.45)
    )
    cap_p = cap_tb.text_frame.paragraphs[0]
    cap_p.alignment = PP_ALIGN.CENTER
    r = cap_p.add_run()
    r.text = "Auto-rotate"
    r.font.size = Pt(14)
    r.font.italic = True
    r.font.color.rgb = MUTED

    # Speaker notes (~30 s) ---------------------------------------------
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = (
        "Under the hood, every chore state change happens in one nightly "
        "Cloud Function — the apps only read. It checks each chore against "
        "its anchor day, uncrosses it if it's due, and rotates the assignee "
        "to the next housemate. It's idempotent, so a missed night just "
        "catches up on the next run. One guard: the very first occurrence "
        "sticks with the seeded user, so a chore you create today won't "
        "silently rotate away before you see it."
    )

    out = Path(__file__).resolve().parent / "chores.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"wrote {path}")
